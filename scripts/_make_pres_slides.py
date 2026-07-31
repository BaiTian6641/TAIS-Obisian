"""生成 reports/TAIS_Obsidian_Final_Pres.pptx（一次性生成器）。

风格对齐 runs/academic_reports_format/Bai Final Pres (1).pptx：
16:9、陈述句式标题、每页一个论点 + 要点 bullet + IEEE 编号引用，末尾 Question/Thank/Reference。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
DARK = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x0F, 0x62, 0xFE)  # Carbon blue
GREY = RGBColor(0x52, 0x56, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, bullets, note=None, title_size=26):
    """title: 陈述句式标题；bullets: [(level, text), ...]；note: 底部小字（引用/出处）。"""
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = DARK

    body = s.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(11.9), Inches(4.9))
    bf = body.text_frame
    bf.word_wrap = True
    for i, (lvl, txt) in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = txt
        p.level = lvl
        p.font.size = Pt(20 if lvl == 0 else 17)
        p.font.color.rgb = DARK if lvl == 0 else GREY
        p.space_after = Pt(8)
    if note:
        nb = s.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5))
        np_ = nb.text_frame.paragraphs[0]
        np_.text = note
        np_.font.size = Pt(11)
        np_.font.color.rgb = GREY
    return s


# 1 封面
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Developing a Self-Learning Edge Language Model"
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = DARK
p = tf.add_paragraph(); p.text = "A Weight Virtual-Memory Architecture and Its 0.1B Pilot Validation"
p.font.size = Pt(22); p.font.color.rgb = ACCENT
p = tf.add_paragraph(); p.text = "Tianrui Bai  |  InterEGR 397  |  University of Wisconsin-Madison"
p.font.size = Pt(16); p.font.color.rgb = GREY

# 2 Objectives
add_slide(
    "The main objectives of this project are to let deployed models keep learning, safely and audibly.",
    [
        (0, "1. Scale the validated 0.1B pilot to a 1B self-learning research model"),
        (1, "10B-token pre-training + 1B-token mid-training annealing (OLMo / SmolLM2 recipes)"),
        (0, "2. Migrate and re-validate endogenous components at 1B"),
        (1, "KAL metacognition, HRL retrieval, HCA injection recall, honest degradation"),
        (0, "3. Extend native context from 1K toward 256K (YaRN + progressive curriculum)"),
        (0, "4. Package for HuggingFace distribution with reproducible inference"),
    ],
    note="Recipes: OLMo 3 Dolmino mid-training [20]; SmolLM2 multi-stage WSD [21]",
)

# 3 问题1：权重冻结
add_slide(
    "Pre-trained LLMs store all knowledge in frozen weights and cannot learn after deployment.",
    [
        (0, "New facts, corrections, and user domains all require expensive re-training"),
        (0, "Mainstream workaround: Retrieval-Augmented Generation (RAG) [1]"),
        (1, "Patches knowledge into the prompt as plain text"),
        (1, "Text is not the model's native knowledge representation"),
        (1, "No persistent memory — the same fact must be re-retrieved every conversation"),
        (0, "Active variants (FLARE [2], Self-RAG [3]) improve timing, not the substrate"),
    ],
    note="[1] Lewis et al., NeurIPS 2020  [2] Jiang et al., EMNLP 2023  [3] Asai et al., ICLR 2024",
)

# 4 问题2：无元认知
add_slide(
    "Models have no readout of their own knowledge boundary — so they hallucinate instead of declining.",
    [
        (0, "Linear probes decode \"the model knows it does not know\" from hidden states"),
        (1, "SAPLMA: internal states signal lying [4]"),
        (1, "0.904–1.000 AUROC even under 4-bit quantization [5]"),
        (0, "The signal exists — but mainstream architectures never use it to drive behavior"),
        (0, "Result: fluent guessing where the model should say \"I don't know\""),
    ],
    note="[4] Azaria & Mitchell, EMNLP Findings 2023  [5] arXiv:2606.02628, 2026",
)

# 5 问题3：长上下文成本
add_slide(
    "Long-context inference is expensive: attention scales quadratically and the KV cache grows linearly.",
    [
        (0, "KV cache grows with context length × model size; one long sequence can take GB of memory [6]"),
        (0, "Edge devices cannot afford either the compute or the memory"),
        (0, "Meanwhile, small/edge LLMs are the fastest-growing deployment segment"),
        (1, "SLM market: USD 0.93B (2025) → 5.45B (2032), CAGR 28.7% [7]"),
        (1, "Drivers: privacy, latency, energy efficiency, data control [8]"),
    ],
    note="[6] arXiv:2603.20397  [7] MarketsandMarkets SLM Report 2025  [8] Kristiani et al., AI 7(1):15, 2026",
)

# 6 问题4：持续学习破坏
add_slide(
    "Naive continual learning is destructive: new data overwrites old weights (catastrophic forgetting).",
    [
        (0, "Even optimizer choice changes forgetting: pre-training's own optimizer forgets less [9]"),
        (0, "The brain solves the same problem with fast + slow complementary learning systems"),
        (1, "Hippocampus: fast episodic capture; neocortex: slow consolidation (CLS theory [10])"),
        (0, "Current LLM architectures have no equivalent mechanism"),
    ],
    note="[9] arXiv:2605.06654  [10] McClelland, McNaughton & O'Reilly, Psych. Review 1995",
)

# 7 核心思想
add_slide(
    "Core idea: treat knowledge as a runtime object at the same level as weights — managed like virtual memory.",
    [
        (0, "KnowledgeBlock: the \"page\" of knowledge (markdown source + compiled form)"),
        (0, "Page table (SQLite) + tiered storage: L0 VRAM ↔ L1 DRAM ↔ L2 NVMe ↔ L3 remote"),
        (0, "Fail-closed page faults: model declares \"memory unavailable\" instead of guessing"),
        (0, "Read/write asymmetry"),
        (1, "Runtime: zero-gradient fast writes only (logs, vectors, KV / memory-layer delta)"),
        (1, "Sleep phase: gated gradient consolidation into the backbone (fast/slow CLS [10])"),
        (0, "Tamper-evident signatures; markdown source is the audit and rollback basis"),
    ],
    note="Design: weight virtual memory, block spec, storage tiers, read/write asymmetry",
)

# 8 主干
add_slide(
    "The hybrid backbone keeps long-context cost near-linear, a prerequisite for edge deployment.",
    [
        (0, "GDN-2 linear attention: constant-size recurrent state, decoupled erase/write gates [13]"),
        (0, "Three-level retrieval attention stack [14]"),
        (1, "L0: 512-token sliding window (exact local attention)"),
        (1, "L1: CSA compressed selection (stride-4)"),
        (1, "L2: HCA heavily compressed gist (128:1)"),
        (0, "PM-stream: 5-stream manifold-constrained hyper-connections [18]"),
    ],
    note="[13] Gated DeltaNet-2, arXiv:2605.22791  [14] NSA, arXiv:2502.11089  [18] mHC, arXiv:2512.24880",
)

# 9 KAL
add_slide(
    "KAL metacognition turns the internal \"I don't know\" signal into behavior, not just monitoring.",
    [
        (0, "Small frozen linear heads read hidden states (read-only, separate layers from intervention)"),
        (0, "Three-state P(IK): knowing / uncertain / unknown + affect + conflict"),
        (0, "Low P(IK) drives action: retrieve, ask, call tool, or honestly decline"),
        (0, "Truth-anchor calibration: anchored on factual truth, not LM confidence"),
    ],
    note="Probe evidence: SAPLMA [4]; quantized probes [5]",
)

# 10 主动求知闭环
add_slide(
    "The active inquiry loop acquires knowledge with verification — never blind self-correction.",
    [
        (0, "LLMs fail at unaided self-correction [15] → every write passes a cross-verification gate"),
        (1, "Multi-source consistency, prior consistency, conflict detection"),
        (0, "Write-then-use: HRL indexer retrieves → HCA injection works in the same conversation"),
        (0, "Sleep consolidation: CA1-style gate + ternary-reward RL [16]"),
        (0, "Dynamic vocabulary: reserved concept slots mint new tokens from the inner lexicon [17]"),
    ],
    note="[15] Huang et al., ICLR 2024  [16] TruthRL, arXiv:2509.25760  [17] Kaplan et al., ICLR 2025",
)

# 11 验证总览
add_slide(
    "Every subsystem was implemented and measured on a 0.1B pilot; 437 unit tests pass.",
    [
        (0, "Backbone ablations: hybrid 3.768 | +tri-stack 3.762 | +PM-stream 3.744 | combined 3.743"),
        (0, "GDN-2 gate convergence: 3-stage evidence chain; bounded decay = 4× faster convergence"),
        (0, "KAL: probe AUROC 0.945; truth-anchor calibration 0.845 / 0.829 (two protocols, 3 seeds)"),
        (0, "HRL block retrieval top-1 = 1.000; HCA injection recall = 0.625 (in-context bound 0.70)"),
        (0, "Honest degradation on fabricated facts: 16/16 declines"),
        (0, "Internalization gap 0.015 → 0.758; dissociation 1.000; sleep gate: 8 promote / 1 quarantine / 8 reject"),
    ],
    note="All numbers from project evaluation artifacts (runs/*/report.json, training logs)",
)

# 12 设计历程
add_slide(
    "Engineering honesty changed the design: four findings that mattered.",
    [
        (0, "GDN-2's early retrieval lag was under-converged gates, not an architecture defect → bounded decay"),
        (0, "Injection recall had a gating side effect (0.688 → 0.250)"),
        (1, "Full decoupling failed honestly → memory-layer relocation restored 0.688 with zero interference"),
        (0, "Logprob is not truth: confidence-anchored calibration plateaued (0.769) → truth-anchor fixed it"),
        (0, "Silent Muon×WSD scheduler bug caught in review, fixed, and locked with a regression test"),
    ],
    note="Negative results are recorded, not hidden (fully-decoupled gate; prediction-feedback loop)",
)

# 13 当前阶段
add_slide(
    "The project has transitioned to a 1B model with a complete, tested training-to-upload toolchain.",
    [
        (0, "1B configuration verified: 1,017.7M params, d_model 1536, 32 layers (8×{3 GDN-2 + 1 stack})"),
        (0, "10B-token corpus: web education 73% / math 12% / synthetic textbooks 10% / Chinese 5%"),
        (0, "Mid-training annealing: 1B tokens, quality-upweighted mixture, lr decays to zero"),
        (0, "Muon optimizer end-to-end (converges better than AdamW: 6.523 vs 6.868)"),
        (0, "Honestly positioned: 10B tokens is an architecture-validation budget"),
        (1, "Half of Chinchilla-optimal [22]; far below current 1B practice (4T+ [21])"),
    ],
    note="[20] OLMo 3, arXiv:2512.13961  [21] SmolLM2, arXiv:2502.02737  [22] Chinchilla, arXiv:2203.15556",
)

# 14 影响与下一步
add_slide(
    "A self-learning edge LLM turns a frozen snapshot into an auditable, personal, accumulating asset.",
    [
        (0, "Privacy: knowledge stays on device; blocks are signed, revocable, and auditable"),
        (0, "Next steps"),
        (1, "Complete 1B training; first observation = KAL probe strength at scale"),
        (1, "Train the memory-layer readout interface (side-effect-free recall)"),
        (1, "256K progressive context curriculum (YaRN, staged)"),
        (1, "Standard-format export for community benchmarking"),
        (0, "Long-term: 1.5B model with native 1M context and the full self-learning loop"),
    ],
    note="Target: edge-native continually-learning LLM",
)

# 15 Question
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(5.2), Inches(3.1), Inches(3), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = "Question?"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = DARK

# 16 Thank
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(5.2), Inches(3.1), Inches(3), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = "Thank you"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = DARK

# 17 Reference
refs = [
    "[1] Lewis et al., \"RAG for knowledge-intensive NLP tasks,\" NeurIPS 2020, arXiv:2005.11401",
    "[2] Jiang et al., \"FLARE,\" EMNLP 2023, arXiv:2305.06983",
    "[3] Asai et al., \"Self-RAG,\" ICLR 2024, arXiv:2310.11511",
    "[4] Azaria & Mitchell, \"SAPLMA,\" EMNLP Findings 2023, arXiv:2304.13734",
    "[5] \"Hallucination linearly decodable in quantized LLMs,\" arXiv:2606.02628, 2026",
    "[6] \"KV cache optimization strategies,\" arXiv:2603.20397, 2026",
    "[7] MarketsandMarkets, \"Small language model market report 2025–2032\"",
    "[8] Kristiani et al., \"Deploying LLM on edge devices: survey,\" AI 7(1):15, 2026",
    "[9] Liu, \"Optimizer-model consistency,\" arXiv:2605.06654, 2026",
    "[10] McClelland et al., \"Complementary learning systems,\" Psych. Review 102(3), 1995",
    "[13] Hatamizadeh et al., \"Gated DeltaNet-2,\" arXiv:2605.22791, 2026",
    "[14] Yuan et al., \"Native sparse attention,\" arXiv:2502.11089, 2025",
    "[15] Huang et al., \"LLMs cannot self-correct reasoning yet,\" ICLR 2024, arXiv:2310.01798",
    "[16] \"TruthRL,\" arXiv:2509.25760, 2025",
    "[17] Kaplan et al., \"From tokens to words,\" ICLR 2025, arXiv:2410.05864",
    "[18] Xie et al., \"mHC: Manifold-constrained hyper-connections,\" arXiv:2512.24880, 2025",
    "[19] Kimi Team, \"Kimi Linear,\" arXiv:2510.26692, 2025",
    "[20] OLMo Team, \"OLMo 3,\" arXiv:2512.13961, 2025",
    "[21] Ben Allal et al., \"SmolLM2,\" arXiv:2502.02737, 2025",
    "[22] Hoffmann et al., \"Chinchilla,\" NeurIPS 2022, arXiv:2203.15556",
]
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8))
p = tb.text_frame.paragraphs[0]
p.text = "Reference"
p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = DARK
body = s.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(11.9), Inches(6.1))
bf = body.text_frame
bf.word_wrap = True
for i, r in enumerate(refs):
    p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
    p.text = r
    p.font.size = Pt(11)
    p.font.color.rgb = GREY
    p.space_after = Pt(4)

out = ROOT / "reports" / "TAIS_Obsidian_Final_Pres.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"written {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
