"""生成 reports/TAIS_Obsidian_汇报幻灯片_中文.pptx（中文版，嵌入架构图与数据图表）。

风格对齐参考样例：16:9、陈述句式标题、每页一个论点 + 要点 bullet + IEEE 编号引用。
图表来源：reports/assets/（_make_report_assets.py 生成，全部真实数据）。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "reports" / "assets"
DARK = RGBColor(0x16, 0x16, 0x16)
ACCENT = RGBColor(0x0F, 0x62, 0xFE)
GREY = RGBColor(0x52, 0x52, 0x52)
TEAL = RGBColor(0x00, 0x9D, 0x9A)
GREEN = RGBColor(0x24, 0xA1, 0x48)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
FONT = "微软雅黑"


def _set(p, size, color, bold=False):
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT


def add_slide(title, bullets, note=None, title_size=25):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], title_size, DARK, True)
    tf.paragraphs[0].text = title
    body = s.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.9), Inches(5.0))
    bf = body.text_frame
    bf.word_wrap = True
    for i, (lvl, txt) in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = txt
        p.level = lvl
        _set(p, 19 if lvl == 0 else 16, DARK if lvl == 0 else GREY)
        p.space_after = Pt(8)
    if note:
        nb = s.shapes.add_textbox(Inches(0.55), Inches(6.9), Inches(12.2), Inches(0.45))
        p = nb.text_frame.paragraphs[0]
        p.text = note
        _set(p, 10.5, GREY)
    return s


def add_image_slide(title, img, note=None, img_h=5.3):
    s = add_slide(title, [], note)
    pic = s.shapes.add_picture(str(ASSETS / img), Inches(0), Inches(1.7), height=Inches(img_h))
    pic.left = int((prs.slide_width - pic.width) / 2)
    return s


# 1 封面
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
_set(tf.paragraphs[0], 38, DARK, True)
tf.paragraphs[0].text = "自学习边缘语言模型"
p = tf.add_paragraph(); _set(p, 21, ACCENT)
p.text = "权重虚拟内存架构及其 0.1B 先导验证 —— 阶段性汇报"
p = tf.add_paragraph(); _set(p, 15, GREY)
p.text = "Tianrui Bai  |  2026 年 7 月"

# 2 目标
add_slide(
    "本项目的目标是让已部署模型持续学习——安全、可审计、面向边缘设备。",
    [
        (0, "1. 把已验证的 0.1B pilot 扩展到 1B 自学习研究模型"),
        (1, "10B tokens 预训练 + 1B tokens 中训练退火（OLMo / SmolLM2 配方）"),
        (0, "2. 在 1B 规模迁移并复测内生部件"),
        (1, "KAL 元认知、HRL 检索、HCA 注入召回、诚实降级"),
        (0, "3. 原生上下文从 1K 扩展到 256K（YaRN + 渐进扩窗）"),
        (0, "4. 打包发布 HuggingFace，保证推理链路可复现"),
    ],
    note="配方依据：OLMo 3 Dolmino 中训练 [20]；SmolLM2 多阶段 WSD [21]",
)

# 3 问题1
add_slide(
    "预训练 LLM 把全部知识存在冻结权重里，部署后无法学习。",
    [
        (0, "新事实、纠错、用户领域适配都要付出昂贵重训练代价"),
        (0, "主流绕行：检索增强生成（RAG）[1]"),
        (1, "把知识以纯文本补丁进提示词"),
        (1, "文本不是模型的原生知识表示"),
        (1, "无持久记忆——同一事实每段对话都要重新检索"),
        (0, "主动变体（FLARE [2]、Self-RAG [3]）改进时机，不改载体"),
    ],
    note="[1] Lewis et al., NeurIPS 2020  [2] Jiang et al., EMNLP 2023  [3] Asai et al., ICLR 2024",
)

# 4 问题2
add_slide(
    "模型读不出自己的知识边界——所以在该拒答的地方产生幻觉。",
    [
        (0, "线性探针可从隐状态解码\"模型知道自己不知道\""),
        (1, "SAPLMA：内部状态暴露说谎 [4]"),
        (1, "4-bit 量化下仍达 0.904–1.000 AUROC [5]"),
        (0, "信号就在那里——但主流架构从不拿它驱动行为"),
        (0, "结果：本该说\"我不知道\"的地方，模型流畅地猜测"),
    ],
    note="[4] Azaria & Mitchell, EMNLP Findings 2023  [5] arXiv:2606.02628, 2026",
)

# 5 问题3
add_slide(
    "长上下文很贵：注意力平方增长，KV 缓存线性增长——边缘设备两头都扛不住。",
    [
        (0, "KV 缓存随上下文长度 × 模型规模增长，单条长序列可达数 GB [6]"),
        (0, "而小/边缘模型恰是增长最快的部署方向"),
        (1, "SLM 市场：9.3 亿美元（2025）→ 54.5 亿（2032），CAGR 28.7% [7]"),
        (1, "驱动力：隐私、延迟、能效、数据控制 [8]"),
        (0, "朴素持续微调还会灾难性遗忘——连优化器选择都影响遗忘量 [9]"),
    ],
    note="[6] arXiv:2603.20397  [7] MarketsandMarkets 2025  [8] Kristiani et al., 2026  [9] arXiv:2605.06654",
)

# 6 核心思想
add_slide(
    "核心思想：把知识升格为与权重同级的运行时对象，用操作系统式虚拟内存管理。",
    [
        (0, "知识块（KnowledgeBlock）= 知识的\"页\"（markdown 源代码 + 编译产物）"),
        (0, "页表（SQLite）+ 分层存储：L0 VRAM ↔ L1 DRAM ↔ L2 NVMe ↔ L3 远端"),
        (0, "缺页 fail-closed：声明\"记忆暂不可用\"，绝不空白作答"),
        (0, "读写不对称"),
        (1, "运行时仅零梯度快写（日志 / 向量 / KV / 记忆层 delta）"),
        (1, "睡眠期才做受门控的梯度固化（快慢互补学习系统 [10]）"),
    ],
    note="[10] McClelland, McNaughton & O'Reilly, Psych. Review 1995",
)

# 7 架构图
add_image_slide(
    "TAIS Obsidian 总体架构：主干、TAIS 内核、主动求知闭环、知识块库、睡眠固化。",
    "architecture_v3.png",
    note="v3.0（IBM Carbon 设计语言）｜主干消融 −0.025 nats；部件全部 checkpoint 内生",
    img_h=5.35,
)

# 8 KAL
add_slide(
    "KAL 元认知把内部的\"我不知道\"信号从监测变成行为。",
    [
        (0, "小型冻结线性头只读隐状态（与干预点分置不同层，防自激）"),
        (0, "三态 P(IK)：知道 / 不确定 / 不知道 + 情感 + 冲突"),
        (0, "低 P(IK) 驱动行为：检索、提问、调工具、或诚实降级"),
        (0, "真值锚校准：锚事实真假，而非语言建模置信度"),
        (0, "主动求知闭环：交叉验证才写入——绝不裸自我修正 [15]"),
    ],
    note="[15] Huang et al., ICLR 2024（LLM 无法独立完成推理自纠）",
)

# 9 训练曲线
add_image_slide(
    "0.1B pilot 预训练 10k 步稳定收敛（真实训练日志）。",
    "chart_training_curve.png",
    note="runs/_gdn2_10k_train_out.txt（每 50 步采样 + EMA 平滑）",
)

# 10 消融
add_image_slide(
    "原生部件消融：PM-stream −0.024、组合 −0.025 nats，代价可控。",
    "chart_ablation.png",
    note="0.1B / 2000 步 val loss；hybrid 基线 3.768",
)

# 11 KAL 校准
add_image_slide(
    "KAL 真值校准从 0.769 提到 0.845 / 0.829 双口径，达成 ≥0.8 目标。",
    "chart_kal.png",
    note="3 seed 均值±std；预测反馈循环无增益已回滚（诚实负结果）",
)

# 12 全链强度
add_image_slide(
    "统一 checkpoint 全链已训强度：检索 0.938 → 召回 0.625 → 内化 0.688 → 诚实降级 16/16。",
    "chart_fullchain.png",
    note="n=16 实测；红色虚线为各基线（线性门控 0.188 / 门控副作用 0.250）",
)

# 13 设计历程
add_slide(
    "工程诚实改变了设计：四个关键发现。",
    [
        (0, "GDN-2 早期检索落后 = 门欠收敛（非架构缺陷）→ 有界 decay 4× 加速"),
        (0, "注入召回的门控副作用（0.688→0.250）"),
        (1, "彻底解耦失败（诚实负结果）→ 记忆层条目迁移根治，零干扰恢复 0.688"),
        (0, "logprob ≠ 真相：置信度锚停在 0.769 → 真值锚 + 锚集扩充达标"),
        (0, "Muon×WSD 静默 bug 审阅捕获、修复、回归测试锁定"),
    ],
    note="负结果均如实记录（彻底解耦门控 / 预测反馈循环）",
)

# 14 当前阶段
add_slide(
    "项目已过渡到 1B：配置、数据、训练到上传的全工具链均已验证就绪。",
    [
        (0, "1B 配置实测：1,017.7M 参数，d1536 × 32 层（8×{3 GDN-2 + 1 三级栈}）"),
        (0, "10B tokens：教育网页 73% / 数学 12% / 合成教科书 10% / 中文 5%"),
        (0, "中训练退火 1B tokens（高质量上移，lr 线性到 0）"),
        (0, "Muon 优化器贯穿（收敛优于 AdamW：6.523 vs 6.868）"),
        (0, "诚实定位：10B 是架构验证预算（Chinchilla 一半 [22]，远低于当代 4T+ [21]）"),
    ],
    note="[20] OLMo 3  [21] SmolLM2  [22] Chinchilla；437 项 pytest 全绿",
)

# 15 NIAH
add_image_slide(
    "NIAH 长度扫描：GDN-2 短中长略优；max_seq=1024 硬限催生 256K 扩容工程。",
    "chart_niah.png",
    note="50 queries/cell；>1024 截断为 0 → RoPE 扩容 + YaRN 渐进课程",
)

# 16 影响
add_slide(
    "自学习边缘 LLM：把冻结快照变成可审计、可撤销、持续累积的个人资产。",
    [
        (0, "隐私内生：知识不出设备；块带签名、可回滚、可审计"),
        (0, "下一步"),
        (1, "完成 1B 训练；首个观测点 = 1B 下 KAL 探针强度"),
        (1, "记忆层读出接口训练（零副作用召回）"),
        (1, "256K 渐进扩窗课程（YaRN，分阶段）"),
        (0, "远期：1.5B 原生 1M 上下文 + 端到端自学习闭环"),
    ],
    note="目标：边缘原生的持续学习 LLM",
)

# 17 Question
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(4.9), Inches(3.1), Inches(4), Inches(1.2))
_set(tb.text_frame.paragraphs[0], 44, DARK, True)
tb.text_frame.paragraphs[0].text = "提问？"

# 18 Thank
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(4.9), Inches(3.1), Inches(4), Inches(1.2))
_set(tb.text_frame.paragraphs[0], 44, DARK, True)
tb.text_frame.paragraphs[0].text = "谢谢"

# 19 Reference
refs = [
    "[1] Lewis et al., \"RAG for knowledge-intensive NLP tasks,\" NeurIPS 2020, arXiv:2005.11401",
    "[2] Jiang et al., \"FLARE,\" EMNLP 2023, arXiv:2305.06983",
    "[3] Asai et al., \"Self-RAG,\" ICLR 2024, arXiv:2310.11511",
    "[4] Azaria & Mitchell, \"SAPLMA,\" EMNLP Findings 2023, arXiv:2304.13734",
    "[5] \"Hallucination linearly decodable in quantized LLMs,\" arXiv:2606.02628, 2026",
    "[6] \"KV cache optimization strategies,\" arXiv:2603.20397, 2026",
    "[7] MarketsandMarkets, \"Small language model market report 2025–2032\"",
    "[8] Kristiani et al., \"Deploying LLM on edge devices: survey,\" AI 7(1):15, 2026",
    "[9] Liu, \"Optimizer-model consistency,\" arXiv:2605.06654, 2026",
    "[10] McClelland et al., \"Complementary learning systems,\" Psych. Review 102(3), 1995",
    "[15] Huang et al., \"LLMs cannot self-correct reasoning yet,\" ICLR 2024, arXiv:2310.01798",
    "[20] OLMo Team, \"OLMo 3,\" arXiv:2512.13961, 2025",
    "[21] Ben Allal et al., \"SmolLM2,\" arXiv:2502.02737, 2025",
    "[22] Hoffmann et al., \"Chinchilla,\" NeurIPS 2022, arXiv:2203.15556",
]
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8))
_set(tb.text_frame.paragraphs[0], 26, DARK, True)
tb.text_frame.paragraphs[0].text = "参考文献（Reference）"
body = s.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(11.9), Inches(6.1))
bf = body.text_frame
bf.word_wrap = True
for i, r in enumerate(refs):
    p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
    p.text = r
    _set(p, 11.5, GREY)
    p.space_after = Pt(4)

out = ROOT / "reports" / "TAIS_Obsidian_汇报幻灯片_中文.pptx"
prs.save(out)
print(f"written {out} ({len(prs.slides._sldIdLst)} slides)")
